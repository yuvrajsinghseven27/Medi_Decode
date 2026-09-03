import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Paths
ROOT_DIR = Path(__file__).resolve().parent.parent.parent
BRAIN_DIR = Path(r"C:\Users\yuvra\.gemini\antigravity\brain\73f8b424-3826-45dd-90dc-bb55744536ae")
PPT_OUT = ROOT_DIR / "MediDecode_Project_Presentation.pptx"

# Color Palette
COLOR_BG = RGBColor(11, 19, 43)        # #0B132B Dark Navy
COLOR_CARD_BG = RGBColor(15, 23, 42)   # #0F172A Slate Card
COLOR_CARD_BORDER = RGBColor(30, 41, 59)
COLOR_TEAL = RGBColor(13, 148, 136)     # #0D9488 Teal
COLOR_CYAN = RGBColor(45, 212, 191)     # #2DD4BF Cyan/Aqua
COLOR_EMERALD = RGBColor(52, 211, 153)  # #34D399 Emerald
COLOR_ROSE = RGBColor(244, 63, 94)      # #F43F5E Rose
COLOR_AMBER = RGBColor(245, 158, 11)    # #F59E0B Amber
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_SLATE_LIGHT = RGBColor(226, 232, 240)
COLOR_SLATE_MUTED = RGBColor(148, 163, 184)


def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG


def add_header(slide, title_text, category_text="MEDIDECODE • CLINICAL INTELLIGENCE PLATFORM"):
    # Category badge
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.3))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_CYAN

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.5), Inches(0.7))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_WHITE


def add_card(slide, left, top, width, height, bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    return shape


def build_presentation():
    prs = Presentation()
    # 16:9 Widescreen layout (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # =========================================================================
    # SLIDE 1: Title Slide
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)

    # Accent Card Backdrop
    card1 = add_card(slide1, Inches(1.0), Inches(1.2), Inches(11.333), Inches(5.1), bg_color=RGBColor(15, 23, 42), border_color=COLOR_TEAL)

    # Brand Title
    tbox = slide1.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.3), Inches(3.8))
    tf = tbox.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "MEDIDECODE CLINICAL SAFETY PLATFORM"
    p0.font.size = Pt(12)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_CYAN

    p1 = tf.add_paragraph()
    p1.text = "MediDecode"
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE
    p1.space_after = Pt(10)

    p2 = tf.add_paragraph()
    p2.text = "AI-Powered Multimodal Prescription Parsing, Cross-Doctor Multi-Script Reconciliation & Cultural Dietary Safety Platform"
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLOR_CYAN
    p2.space_after = Pt(20)

    p3 = tf.add_paragraph()
    p3.text = "• Google Gemini 2.5 / 3.6 Flash Multimodal Vision • Cross-Doctor Polypharmacy De-Duplication\n• Localized Cultural Dietary Rules (Chai, Dairy, Fasting) • Chrono-Medication Scheduling & Refill Tracker"
    p3.font.size = Pt(13)
    p3.font.color.rgb = COLOR_SLATE_LIGHT
    p3.space_after = Pt(24)

    p4 = tf.add_paragraph()
    p4.text = "Repository: yuvrajsinghseven27/Medi_Decode • Production Validated (18/18 Tests Passing)"
    p4.font.size = Pt(11)
    p4.font.color.rgb = COLOR_SLATE_MUTED

    # =========================================================================
    # SLIDE 2: The Healthcare Problem
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "The Critical Problem: Fragmented Care & Silent Overdoses")

    # 4 Problem Cards in a Grid
    problems = [
        ("Doctor Silos & Polypharmacy", "Patients consult multiple independent specialists (Cardiologist, Orthopedist, GP) who have no shared electronic record, creating fragmented drug regimens.", COLOR_ROSE),
        ("Duplicate Brand Disguise", "Identical active molecules are prescribed under disparate commercial brand names (e.g., Dolo 650 + Crocin Pain Relief = both Paracetamol), hiding lethal overlaps.", COLOR_ROSE),
        ("Unchecked Cumulative Toxicity", "Patients take multiple tablets daily, innocently exceeding safe clinical ceilings (e.g. Paracetamol > 4,000mg/day) leading to acute liver necrosis and hospital readmission.", COLOR_AMBER),
        ("Unaddressed Cultural Diets", "Traditional diets severely interfere with pharmacokinetics: high-tannin milk tea (chai) blunts thyroid/iron absorption; dairy chelates antibiotics; pickles counteract blood pressure meds.", COLOR_AMBER),
    ]

    for idx, (title, desc, accent) in enumerate(problems):
        col = idx % 2
        row = idx // 2
        left = Inches(1.0 + col * 5.8)
        top = Inches(1.6 + row * 2.6)
        add_card(slide2, left, top, Inches(5.5), Inches(2.3), border_color=accent)

        tb = slide2.shapes.add_textbox(left + Inches(0.3), top + Inches(0.2), Inches(4.9), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True
        p_t = tf.paragraphs[0]
        p_t.text = f"⚠️ {title}"
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = accent
        p_t.space_after = Pt(8)

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = COLOR_SLATE_LIGHT

    # =========================================================================
    # SLIDE 3: MediDecode 4 Core Pillars Solution
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "The MediDecode Solution: 4 Clinical Intelligence Pillars")

    pillars = [
        ("1. Multimodal Vision OCR", "Ingests doctor handwriting using Google Gemini 2.5/3.6 Flash. Emits deterministic JSON via Pydantic response schema with per-field confidence scoring.", COLOR_CYAN),
        ("2. Cross-Doctor Reconciliation", "Standardizes brands to active generic molecules. Calculates cumulative daily mg intake, enforces safe ceilings, and generates exportable Doctor Consultation Notes.", COLOR_ROSE),
        ("3. Cultural Dietary & Fasting Guard", "Engineered for Indian diets: warns against milk tea (chai) tannin blunting, dairy chelation, and dynamically adapts dosages for Ramadan (Suhoor/Iftar) and Navratri.", COLOR_AMBER),
        ("4. Chrono-Scheduling & Refill", "Anchors dosing to patient-specific meal times (AC -30m, PC +30m). Decrements pill inventory on intake and warns when medication supply drops ≤ 3 days.", COLOR_EMERALD),
    ]

    for idx, (title, desc, accent) in enumerate(pillars):
        left = Inches(0.8 + idx * 2.95)
        top = Inches(1.6)
        add_card(slide3, left, top, Inches(2.8), Inches(5.1), border_color=accent)

        tb = slide3.shapes.add_textbox(left + Inches(0.2), top + Inches(0.3), Inches(2.4), Inches(4.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = accent
        p_t.space_after = Pt(14)

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = COLOR_SLATE_LIGHT
        p_d.space_after = Pt(14)

    # =========================================================================
    # SLIDE 4: Multimodal Vision & Split-Screen Verification
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "Multimodal Document Vision & Split-Screen Verification")

    # Left: Explanation & Architecture
    add_card(slide4, Inches(0.8), Inches(1.6), Inches(4.6), Inches(5.2), border_color=COLOR_CYAN)
    tb = slide4.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(4.2), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Clinical Handwriting OCR Architecture"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    p.space_after = Pt(10)

    bullets = [
        ("Gemini Multimodal Ingestion", "Upload JPEG, PNG, WEBP, or PDF scans. Uses official google-genai SDK with resilient fallback cascade across 3.6-Flash and Flash-Latest."),
        ("Deterministic Structured Output", "response_schema=PrescriptionExtractionResult enforces 100% typed clinical output: brand, generic, dosage form, strength, frequency, timing relation."),
        ("Low-Confidence Highlighting", "If any medication certainty is < 85% or if cursive marginal notes exist, requires_verification=True flags the card in glowing amber for human confirmation."),
        ("Interactive Editing", "Inline inputs allow doctor or patient to correct dosage or duration before executing safety reconciliation."),
    ]
    for b_title, b_desc in bullets:
        pb = tf.add_paragraph()
        pb.text = f"• {b_title}: {b_desc}"
        pb.font.size = Pt(10.5)
        pb.font.color.rgb = COLOR_SLATE_LIGHT
        pb.space_after = Pt(8)

    # Right: Embedded Screenshot
    img_path = BRAIN_DIR / "split_screen_verification.png"
    if img_path.exists():
        slide4.shapes.add_picture(str(img_path), Inches(5.6), Inches(1.6), width=Inches(6.9))

    # =========================================================================
    # SLIDE 5: Cross-Doctor Reconciliation & Toxicity Ceilings
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "Cross-Doctor De-Duplication & Cumulative Toxicity Engine")

    # Left: Methodology & Ceilings
    add_card(slide5, Inches(0.8), Inches(1.6), Inches(4.6), Inches(5.2), border_color=COLOR_ROSE)
    tb = slide5.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(4.2), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Multi-Script Polypharmacy Defense"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_ROSE
    p.space_after = Pt(10)

    bullets5 = [
        ("Canonical Molecule Normalizer", "Maps Indian & global brand names (Crocin, Dolo, Calpol, Tylenol, Pacimol) to canonical active pharmaceutical ingredients (Paracetamol)."),
        ("Cumulative Daily Intake Math", "Multiplies parsed strength by frequency (OD:1, BD:2, TID:3, QID:4) across all active patient prescriptions to calculate total daily burden."),
        ("Strict Toxic Thresholds", "Paracetamol > 4,000mg/day triggers CRITICAL hepatotoxicity alert. Concomitant NSAIDs (Ibuprofen + Diclofenac) flag compounded gastric ulcer/renal failure risk."),
        ("Doctor Query Summary Sheet", "Compiles printable consultation note detailing overlapping molecules, clinical overages, and discussion prompts for the attending physician."),
    ]
    for b_title, b_desc in bullets5:
        pb = tf.add_paragraph()
        pb.text = f"• {b_title}: {b_desc}"
        pb.font.size = Pt(10.5)
        pb.font.color.rgb = COLOR_SLATE_LIGHT
        pb.space_after = Pt(8)

    # Right: Screenshots (Safety Panel + Doctor Note Modal)
    img_safety = BRAIN_DIR / "safety_panel.png"
    if img_safety.exists():
        slide5.shapes.add_picture(str(img_safety), Inches(5.6), Inches(1.6), width=Inches(6.9), height=Inches(2.5))
    img_doc = BRAIN_DIR / "doctor_query_modal.png"
    if img_doc.exists():
        slide5.shapes.add_picture(str(img_doc), Inches(5.6), Inches(4.3), width=Inches(6.9), height=Inches(2.5))

    # =========================================================================
    # SLIDE 6: Cultural Dietary & Religious Fasting Rules
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, "Cultural Dietary Safety & Multilingual Localization")

    # Left: Rules Matrix
    add_card(slide6, Inches(0.8), Inches(1.6), Inches(5.2), Inches(5.2), border_color=COLOR_AMBER)
    tb = slide6.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(4.8), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Localized Dietary & Fasting Matrix"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_AMBER
    p.space_after = Pt(8)

    rules = [
        ("High-Tannin Milk Tea (Chai)", "Tannins and milk casein bind to Levothyroxine, Iron, and Calcium in the gut, severely blunting absorption. 2-hour separation strictly enforced."),
        ("Dairy & Calcium Products", "Milk, yogurt (dahi), and paneer chelate Fluoroquinolones (Ciprofloxacin) and Tetracyclines (Doxycycline) into non-absorbable complexes."),
        ("High-Sodium Diet (Achaar/Papad)", "Traditional savory sodium intake counteracts blood-pressure lowering efficacy of ACE inhibitors and ARBs."),
        ("Ramadan Fasting Sync", "Re-adjusts morning empty-stomach medications to pre-dawn Suhoor (~04:30 AM) and evening doses to post-sunset Iftar (~07:15 PM)."),
        ("5 Regional Dialects", "Empathetic, colloquial plain-language translations in English, हिन्दी (Hindi), தமிழ் (Tamil), తెలుగు (Telugu), and বাংলা (Bengali)."),
    ]
    for r_title, r_desc in rules:
        pb = tf.add_paragraph()
        pb.text = f"• {r_title}: {r_desc}"
        pb.font.size = Pt(10)
        pb.font.color.rgb = COLOR_SLATE_LIGHT
        pb.space_after = Pt(6)

    # Right: Screenshot Hindi Schedule
    img_hindi = BRAIN_DIR / "schedule_hindi.png"
    if img_hindi.exists():
        slide6.shapes.add_picture(str(img_hindi), Inches(6.2), Inches(1.6), width=Inches(6.3))

    # =========================================================================
    # SLIDE 7: Chrono-Scheduling & Smart Refill Depletion
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_header(slide7, "Personalized Chrono-Scheduling & Inventory Depletion")

    # Left: Explanation
    add_card(slide7, Inches(0.8), Inches(1.6), Inches(4.8), Inches(5.2), border_color=COLOR_EMERALD)
    tb = slide7.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(4.4), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Patient-Anchored Chrono-Timing"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_EMERALD
    p.space_after = Pt(8)

    sched_pts = [
        ("Personal Meal Anchors", "Instead of arbitrary alarm clocks, doses are calculated relative to patient's actual waking (06:00), breakfast (08:00), lunch (13:00), and dinner (20:30)."),
        ("Clinical Timing Relations", "AC doses scheduled 30m before meals; PC doses 30m after meals; WITH_FOOD at meal time; Statins automatically placed at bedtime."),
        ("Smart Refill Depletion", "Tracks remaining pills against daily velocity (OD:1, BD:2, TID:3). Predicts runout and warns when inventory is ≤ 3 days of supply."),
        ("Micro-Confetti Adherence", "Marking doses as 'TAKEN' fires a celebratory particle burst, updates real-time daily adherence gauge, and decrements pill inventory."),
    ]
    for s_title, s_desc in sched_pts:
        pb = tf.add_paragraph()
        pb.text = f"• {s_title}: {s_desc}"
        pb.font.size = Pt(10.5)
        pb.font.color.rgb = COLOR_SLATE_LIGHT
        pb.space_after = Pt(8)

    # Right: Dashboard Screenshot
    img_dash = BRAIN_DIR / "dashboard_desktop.png"
    if img_dash.exists():
        slide7.shapes.add_picture(str(img_dash), Inches(5.8), Inches(1.6), width=Inches(6.7))

    # =========================================================================
    # SLIDE 8: System Architecture, Data Layer & Docker Deployment
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8)
    add_header(slide8, "Production Architecture & Multi-Container Deployment")

    # 3 Service Cards
    services = [
        ("Backend API Gateway", "FastAPI • Python 3.11-slim\n• Uvicorn async server on port 8000\n• Official google-genai SDK\n• Resilient multi-model fallback\n• SQLAlchemy 2.0 Async ORM\n• Healthcheck: GET /healthz", COLOR_CYAN),
        ("Client Web Application", "React 19 • Vite • Tailwind v4\n• Multi-stage Nginx Alpine on port 80\n• SPA client-side routing fallback\n• Live health polling & toast alerts\n• Axios REST API service layer\n• Canvas Confetti reward animation", COLOR_EMERALD),
        ("Relational Database", "PostgreSQL 16 Alpine\n• asyncpg asynchronous connection\n• Cross-compatibility with SQLite\n• Automated Alembic migrations\n• Persistent volume: postgres_data\n• Healthcheck: pg_isready", COLOR_AMBER),
    ]
    for idx, (title, details, accent) in enumerate(services):
        left = Inches(0.8 + idx * 3.95)
        top = Inches(1.6)
        add_card(slide8, left, top, Inches(3.8), Inches(3.2), border_color=accent)

        tb = slide8.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(3.4), Inches(2.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = accent
        p_t.space_after = Pt(8)

        p_d = tf.add_paragraph()
        p_d.text = details
        p_d.font.size = Pt(10.5)
        p_d.font.color.rgb = COLOR_SLATE_LIGHT

    # Bottom Banner: Automated Test Status
    add_card(slide8, Inches(0.8), Inches(5.1), Inches(11.733), Inches(1.7), bg_color=RGBColor(15, 23, 42), border_color=COLOR_TEAL)
    tb_b = slide8.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(11.3), Inches(1.5))
    tf_b = tb_b.text_frame
    tf_b.word_wrap = True

    p_bt = tf_b.paragraphs[0]
    p_bt.text = "Automated Pytest Suite Verification: 18 OF 18 TESTS PASSED (100%)"
    p_bt.font.size = Pt(13)
    p_bt.font.bold = True
    p_bt.font.color.rgb = COLOR_CYAN
    p_bt.space_after = Pt(4)

    p_bd = tf_b.add_paragraph()
    p_bd.text = "• Healthcheck probes (2/2)  • Relational persistence & schemas (2/2)  • Multimodal vision OCR & verification flags (5/5)\n• Cross-doctor duplicate molecule & toxicity engine (6/6)  • Chrono-scheduling & smart refill depletion (3/3)\n• Live Gemini Multimodal Vision API inference validated on scanned JPEG document."
    p_bd.font.size = Pt(10)
    p_bd.font.color.rgb = COLOR_SLATE_LIGHT

    # =========================================================================
    # SLIDE 9: Clinical Impact & Roadmap
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9)
    add_header(slide9, "Clinical Impact, Deployment & Future Roadmap")

    # 3 Summary Cards
    impact_cards = [
        ("Measurable Clinical Impact", [
            "Eliminates silent duplicate molecule overdoses across disparate doctors.",
            "Prevents acute hepatotoxicity from cumulative Paracetamol overages.",
            "Increases drug bioavailability through cultural meal/chai separation.",
            "Democratizes prescription literacy with 5 regional Indian dialects.",
        ], COLOR_CYAN),
        ("Production Deployment Ready", [
            "Single-command launch via root docker-compose.yml.",
            "Automated clinical polypharmacy test seeding script (seed.py).",
            "Zero security leaks: API keys enforced via .gitignore & environment vars.",
            "OpenAPI / Swagger documentation available at /docs.",
        ], COLOR_EMERALD),
        ("Strategic Future Roadmap", [
            "WhatsApp & automated IVR regional voice alerts for elderly patients.",
            "Continuous glucose & blood pressure wearable biometric sync.",
            "National Digital Health Mission (ABDM / FHIR) record interchange.",
            "Pharmacist validation portal for dispensing sign-off.",
        ], COLOR_AMBER),
    ]

    for idx, (title, items, accent) in enumerate(impact_cards):
        left = Inches(0.8 + idx * 3.95)
        top = Inches(1.6)
        add_card(slide9, left, top, Inches(3.8), Inches(4.5), border_color=accent)

        tb = slide9.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(3.4), Inches(4.1))
        tf = tb.text_frame
        tf.word_wrap = True
        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = accent
        p_t.space_after = Pt(12)

        for item in items:
            pb = tf.add_paragraph()
            pb.text = f"✓ {item}"
            pb.font.size = Pt(10)
            pb.font.color.rgb = COLOR_SLATE_LIGHT
            pb.space_after = Pt(8)

    # Save presentation
    prs.save(str(PPT_OUT))
    print(f"SUCCESS! Presentation successfully generated at: {PPT_OUT}")


if __name__ == "__main__":
    build_presentation()
